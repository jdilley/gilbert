"""Text extraction — extract plain text from various document formats."""

import io
import json
import logging
from typing import Any

from gilbert.interfaces.knowledge import DocumentContent, DocumentType

logger = logging.getLogger(__name__)


def extract_text(content: DocumentContent) -> str:
    """Extract plain text from a document based on its type."""
    match content.meta.document_type:
        case DocumentType.TEXT | DocumentType.MARKDOWN | DocumentType.CSV:
            return content.data.decode(content.encoding, errors="replace")
        case DocumentType.JSON:
            return _extract_json(content.data, content.encoding)
        case DocumentType.YAML:
            return _extract_yaml(content.data, content.encoding)
        case DocumentType.PDF:
            return _extract_pdf(content.data)
        case DocumentType.WORD:
            return _extract_word(content.data)
        case DocumentType.EXCEL:
            return _extract_excel(content.data)
        case DocumentType.POWERPOINT:
            return _extract_powerpoint(content.data)
        case _:
            return content.data.decode(content.encoding, errors="replace")


def _extract_json(data: bytes, encoding: str = "utf-8") -> str:
    """Pretty-print JSON for searchability."""
    try:
        parsed = json.loads(data.decode(encoding, errors="replace"))
        return json.dumps(parsed, indent=2, ensure_ascii=False)
    except (json.JSONDecodeError, ValueError):
        return data.decode(encoding, errors="replace")


def _extract_yaml(data: bytes, encoding: str = "utf-8") -> str:
    """Load and dump YAML as formatted text."""
    try:
        import yaml

        parsed = yaml.safe_load(data.decode(encoding, errors="replace"))
        return yaml.dump(parsed, default_flow_style=False, allow_unicode=True)
    except Exception:
        return data.decode(encoding, errors="replace")


def _extract_pdf(data: bytes) -> str:
    """Extract text from PDF using pypdf."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- Page {i} ---\n{text}")
        return "\n\n".join(pages)
    except Exception:
        logger.warning("Failed to extract text from PDF", exc_info=True)
        return ""


def _extract_word(data: bytes) -> str:
    """Extract text from Word documents using python-docx."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        parts: list[str] = []

        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append("\t".join(cells))

        return "\n".join(parts)
    except Exception:
        logger.warning("Failed to extract text from Word document", exc_info=True)
        return ""


def _extract_excel(data: bytes) -> str:
    """Extract text from Excel workbooks using openpyxl."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        parts: list[str] = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append("\t".join(cells))

        wb.close()
        return "\n".join(parts)
    except Exception:
        logger.warning("Failed to extract text from Excel workbook", exc_info=True)
        return ""


def _extract_powerpoint(data: bytes) -> str:
    """Extract text from PowerPoint presentations using python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
            if slide_texts:
                parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))

        return "\n\n".join(parts)
    except Exception:
        logger.warning("Failed to extract text from PowerPoint", exc_info=True)
        return ""
